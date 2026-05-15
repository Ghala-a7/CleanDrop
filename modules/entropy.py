"""
Entropy Analysis Module — CleanDrop
=====================================
WHAT THIS MODULE DOES:
    Measures how random the bytes inside a file are (Shannon Entropy).

    Normal files (text, code) → LOW entropy (patterns exist).
    Encrypted / packed malware → HIGH entropy (looks like noise).

    Three categories of files:
    1. Binary compressed (images, zip, mp4…) — raw entropy is always high
       by design → always SAFE, skip scoring.
    2. Office Open XML (.docx, .xlsx…) — these are ZIP archives internally,
       so raw byte entropy is always ≥7.5 and useless. We extract the text
       content and score entropy on that instead.
    3. Everything else (.txt, .pdf, .exe…) — score raw byte entropy normally.

ENTROPY SCALE (for text/raw content):
    < 5.0   → SAFE     (score:  0)  — plain text, structured data
    5.0–5.9 → LOW      (score: 10)  — slightly above normal
    6.0–6.9 → LOW      (score: 30)  — compressed or mixed content
    7.0–7.1 → MEDIUM   (score: 50)  — suspicious
    7.2–7.4 → HIGH     (score: 65)  — likely obfuscated
    7.5–7.6 → HIGH     (score: 75)  — likely packed/encrypted
    7.7–7.8 → CRITICAL (score: 85)  — strongly encrypted
    >= 7.9  → CRITICAL (score: 95)  — maximum randomness
"""

import math
import os
from collections import Counter

# Binary formats where high entropy is ALWAYS expected — skip scoring
COMPRESSED_EXTENSIONS = {
    '.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp', '.tiff',
    '.zip', '.gz', '.rar', '.7z', '.tar',
    '.mp3', '.mp4', '.avi', '.mov',
}

# Office Open XML formats — ZIP-based internally, score text content instead of raw bytes
OFFICE_EXTENSIONS = {
    '.docx', '.docm',
    '.xlsx', '.xlsm',
    '.pptx', '.pptm',
}


def _byte_entropy(data: bytes) -> float:
    if not data:
        return 0.0
    counter = Counter(data)
    total = len(data)
    return round(-sum((c / total) * math.log2(c / total) for c in counter.values()), 4)


def _score(entropy: float):
    if entropy < 5.0:   return 'SAFE',     0
    elif entropy < 6.0: return 'LOW',      10
    elif entropy < 7.0: return 'LOW',      30
    elif entropy < 7.2: return 'MEDIUM',   50
    elif entropy < 7.5: return 'HIGH',     65
    elif entropy < 7.7: return 'HIGH',     75
    elif entropy < 7.9: return 'CRITICAL', 85
    else:               return 'CRITICAL', 95


def _office_content_entropy(file_path: str, ext: str):
    """
    Extract text from an Office file and return entropy of that content.
    Returns None if extraction fails or yields suspiciously little text,
    which is itself a signal worth flagging as MEDIUM.
    """
    try:
        text = ''

        if ext in ('.docx', '.docm'):
            from docx import Document
            doc = Document(file_path)
            for para in doc.paragraphs:
                text += para.text + ' '
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text += cell.text + ' '

        elif ext in ('.xlsx', '.xlsm'):
            import pandas as pd
            df = pd.read_excel(file_path, dtype=str)
            text = df.to_string()

        elif ext in ('.pptx', '.pptm'):
            try:
                from pptx import Presentation
                prs = Presentation(file_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, 'text'):
                            text += shape.text + ' '
            except ImportError:
                return None

        data = text.encode('utf-8', errors='ignore')
        file_size = os.path.getsize(file_path)

        # Very little extractable text for a non-trivial file size → suspicious
        if file_size > 5_000 and len(data) < file_size * 0.001:
            return None

        return _byte_entropy(data) if data else 0.0

    except Exception:
        return None


def calculate_entropy(file_path: str) -> dict:
    ext = os.path.splitext(file_path)[1].lower()

    try:
        # ── Category 1: Binary compressed formats — always SAFE ──────────────
        if ext in COMPRESSED_EXTENSIONS:
            with open(file_path, 'rb') as f:
                data = f.read()
            entropy = _byte_entropy(data)
            return {
                'entropy': entropy,
                'risk': 'SAFE',
                'score': 0,
                'source': 'Entropy Analysis',
                'note': 'Compressed format — high entropy is expected',
            }

        # ── Category 2: Office Open XML — score text content, not raw bytes ──
        if ext in OFFICE_EXTENSIONS:
            content_entropy = _office_content_entropy(file_path, ext)

            if content_entropy is None:
                # Extraction failed or returned suspiciously little content
                return {
                    'entropy': 0.0,
                    'risk': 'MEDIUM',
                    'score': 50,
                    'source': 'Entropy Analysis',
                    'note': 'Office file — could not extract readable content',
                }

            risk, score = _score(content_entropy)
            return {
                'entropy': content_entropy,
                'risk': risk,
                'score': score,
                'source': 'Entropy Analysis',
                'note': 'Office file — entropy scored on extracted text content',
            }

        # ── Category 3: All other files — raw byte entropy ───────────────────
        with open(file_path, 'rb') as f:
            data = f.read()

        if not data:
            return {'entropy': 0.0, 'risk': 'SAFE', 'score': 0, 'source': 'Entropy Analysis'}

        entropy = _byte_entropy(data)
        risk, score = _score(entropy)

        return {
            'entropy': entropy,
            'risk': risk,
            'score': score,
            'source': 'Entropy Analysis',
        }

    except FileNotFoundError:
        return {'entropy': 0.0, 'risk': 'SAFE', 'score': 0, 'source': 'Entropy Analysis'}


if __name__ == '__main__':
    import sys
    path = sys.argv[1] if len(sys.argv) > 1 else None
    if path and os.path.exists(path):
        print(calculate_entropy(path))
    else:
        print("Usage: python entropy.py <file_path>")
